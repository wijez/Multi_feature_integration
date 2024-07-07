import os
from tkinter import filedialog

import pandas as pd
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation


def read_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    if file_extension.lower() == '.txt':
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    elif file_extension.lower() in ['.doc', '.docx']:
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif file_extension.lower() in ['.xls', '.xlsx']:
        df = pd.read_excel(file_path)
        return df.to_csv(index=False)
    elif file_extension.lower() == '.pdf':
        pdf = PdfReader(file_path)
        text = ''
        for page in pdf.pages:
            text += page.extract_text() + '\n'
        return text
    elif file_extension.lower() == '.pptx':
        presentation = Presentation(file_path)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    else:
        raise ValueError("Unsupported file format")


def write_file(output_file, translated_text):
    with open(output_file, 'w', encoding='utf-8') as file:
        file.write(translated_text)



